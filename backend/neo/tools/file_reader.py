"""File reader — Parse documents into structured text for LLM context.

Supports: Excel (.xlsx/.xls), Word (.docx), PowerPoint (.pptx),
          PDF (.pdf), CSV/TSV, and plain text files.

Multimodal support: Extracts embedded images from PPTX, DOCX, and PDF
files for vision-capable LLMs.
"""

import base64
import csv
import io
import logging
import os
from dataclasses import dataclass, field
from typing import Any

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
MAX_FILE_SIZE = 10 * 1024 * 1024  # 10 MB per file
MAX_TOTAL_CHARS = 100_000  # Max combined output characters

# Image extraction limits
_MAX_IMAGE_DIMENSION = 1568  # Claude's recommended max px
_MAX_IMAGE_BYTES = 2 * 1024 * 1024  # 2 MB per image after resize
_MAX_IMAGES_PER_FILE = 5
_MAX_IMAGES_TOTAL = 10

# MIME types for common image formats
_IMAGE_MIME_MAP = {
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".gif": "image/gif",
    ".webp": "image/webp",
    ".bmp": "image/png",  # convert BMP → PNG
    ".tiff": "image/png",  # convert TIFF → PNG
    ".tif": "image/png",
}


# ---------------------------------------------------------------------------
# Multimodal data structures
# ---------------------------------------------------------------------------
@dataclass
class ImageData:
    """A single extracted image with metadata."""

    data: bytes  # raw image bytes (PNG or JPEG)
    media_type: str  # "image/png" or "image/jpeg"
    source_file: str  # originating file path
    description: str  # e.g. "Slide 3 image" or "Page 2 render"


@dataclass
class ParseResult:
    """Combined text + images from parsing a file."""

    text: str
    images: list[ImageData] = field(default_factory=list)

# Extensions mapped to reader functions
_TEXT_EXTENSIONS = {
    ".txt", ".md", ".json", ".xml", ".log", ".py", ".js", ".ts",
    ".html", ".css", ".yaml", ".yml", ".toml", ".ini", ".cfg",
    ".sh", ".bat", ".sql", ".r", ".java", ".c", ".cpp", ".h",
    ".go", ".rs", ".rb", ".php", ".swift", ".kt", ".scala",
}

# Protected system directories (read-safety)
_BLOCKED_PREFIXES = [
    "/bin", "/sbin", "/usr/bin", "/usr/sbin", "/etc/shadow",
    "/etc/passwd", "/proc", "/sys", "/dev",
]


# ---------------------------------------------------------------------------
# Image helpers (Pillow — lazy import)
# ---------------------------------------------------------------------------
def _resize_image(data: bytes, max_dim: int = _MAX_IMAGE_DIMENSION) -> tuple[bytes, str]:
    """Resize an image if any dimension exceeds max_dim.

    Also converts RGBA → RGB and non-standard formats to PNG.
    Returns (bytes, media_type).
    """
    try:
        from PIL import Image
    except ImportError:
        logger.debug("Pillow not installed — returning raw image bytes")
        return data, "image/png"

    img = Image.open(io.BytesIO(data))

    # Convert palette / RGBA → RGB for JPEG compatibility
    if img.mode in ("RGBA", "LA", "P"):
        background = Image.new("RGB", img.size, (255, 255, 255))
        if img.mode == "P":
            img = img.convert("RGBA")
        background.paste(img, mask=img.split()[-1] if "A" in img.mode else None)
        img = background
    elif img.mode != "RGB":
        img = img.convert("RGB")

    # Resize if needed
    w, h = img.size
    if w > max_dim or h > max_dim:
        ratio = min(max_dim / w, max_dim / h)
        new_size = (int(w * ratio), int(h * ratio))
        img = img.resize(new_size, Image.LANCZOS)

    # Save as JPEG (smaller) for photos, PNG for everything else
    buf = io.BytesIO()
    # Try JPEG first (smaller), fall back to PNG
    img.save(buf, format="JPEG", quality=85)
    result = buf.getvalue()
    media_type = "image/jpeg"

    # If still too large, reduce quality
    if len(result) > _MAX_IMAGE_BYTES:
        buf = io.BytesIO()
        img.save(buf, format="JPEG", quality=60)
        result = buf.getvalue()

    return result, media_type


def _extract_pptx_images(file_path: str) -> list[ImageData]:
    """Extract embedded images from a PowerPoint file."""
    from pptx import Presentation

    images: list[ImageData] = []
    try:
        prs = Presentation(file_path)
    except Exception as e:
        logger.warning("Failed to open PPTX for images: %s", e)
        return images

    for slide_num, slide in enumerate(prs.slides, 1):
        if len(images) >= _MAX_IMAGES_PER_FILE:
            break
        for shape in slide.shapes:
            if len(images) >= _MAX_IMAGES_PER_FILE:
                break
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    blob = shape.image.blob
                    if blob and len(blob) > 100:  # skip tiny placeholder images
                        data, media_type = _resize_image(blob)
                        images.append(ImageData(
                            data=data,
                            media_type=media_type,
                            source_file=os.path.basename(file_path),
                            description=f"Slide {slide_num} image",
                        ))
                except Exception as e:
                    logger.debug("Failed to extract PPTX image: %s", e)

    return images


def _extract_docx_images(file_path: str) -> list[ImageData]:
    """Extract embedded images from a Word document."""
    images: list[ImageData] = []
    try:
        from docx import Document
        doc = Document(file_path)
    except Exception as e:
        logger.warning("Failed to open DOCX for images: %s", e)
        return images

    for rel in doc.part.rels.values():
        if len(images) >= _MAX_IMAGES_PER_FILE:
            break
        if "image" in rel.reltype:
            try:
                blob = rel.target_part.blob
                if blob and len(blob) > 100:
                    data, media_type = _resize_image(blob)
                    images.append(ImageData(
                        data=data,
                        media_type=media_type,
                        source_file=os.path.basename(file_path),
                        description=f"Document image ({rel.target_ref})",
                    ))
            except Exception as e:
                logger.debug("Failed to extract DOCX image: %s", e)

    return images


def _extract_pdf_images(file_path: str, max_pages: int = 5) -> list[ImageData]:
    """Render PDF pages as images via pdfplumber."""
    images: list[ImageData] = []
    try:
        import pdfplumber
    except ImportError:
        logger.debug("pdfplumber not available for PDF image extraction")
        return images

    try:
        with pdfplumber.open(file_path) as pdf:
            pages_to_render = min(len(pdf.pages), max_pages)
            for i in range(pages_to_render):
                if len(images) >= _MAX_IMAGES_PER_FILE:
                    break
                try:
                    page = pdf.pages[i]
                    img = page.to_image(resolution=150)
                    buf = io.BytesIO()
                    img.save(buf, format="PNG")
                    raw = buf.getvalue()
                    data, media_type = _resize_image(raw)
                    images.append(ImageData(
                        data=data,
                        media_type=media_type,
                        source_file=os.path.basename(file_path),
                        description=f"Page {i + 1} render",
                    ))
                except Exception as e:
                    logger.debug("Failed to render PDF page %d: %s", i + 1, e)
    except Exception as e:
        logger.warning("Failed to open PDF for image extraction: %s", e)

    return images


def _describe_excel_charts(file_path: str) -> str:
    """Return text descriptions of chart objects in an Excel file."""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(file_path, data_only=True, read_only=False)
    except Exception:
        return ""

    descriptions: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        charts = getattr(ws, "_charts", [])
        for chart in charts:
            chart_type = type(chart).__name__
            title = getattr(chart, "title", None) or "Untitled"
            descriptions.append(f"[Chart in '{sheet_name}': {chart_type} — {title}]")

    wb.close()
    return "\n".join(descriptions)


# ---------------------------------------------------------------------------
# Security
# ---------------------------------------------------------------------------
def _validate_read_path(file_path: str) -> str:
    """Validate and resolve a file path for reading.

    Returns the resolved absolute path.
    Raises ValueError for blocked paths or missing files.
    """
    expanded = os.path.expanduser(file_path)
    real_path = os.path.realpath(expanded)

    # Block sensitive system paths
    for prefix in _BLOCKED_PREFIXES:
        if real_path.startswith(prefix):
            raise ValueError(f"Cannot read from protected path: {real_path}")

    if not os.path.isfile(real_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    size = os.path.getsize(real_path)
    if size > MAX_FILE_SIZE:
        raise ValueError(
            f"File too large ({size / 1024 / 1024:.1f} MB). "
            f"Maximum is {MAX_FILE_SIZE / 1024 / 1024:.0f} MB."
        )

    return real_path


# ---------------------------------------------------------------------------
# Excel reader
# ---------------------------------------------------------------------------
def read_excel(file_path: str, max_rows: int = 200) -> str:
    """Read an Excel workbook into structured markdown tables.

    Args:
        file_path: Path to .xlsx or .xls file.
        max_rows: Maximum rows to read per sheet.

    Returns:
        Formatted markdown string with sheet data.
    """
    from openpyxl import load_workbook

    path = _validate_read_path(file_path)
    wb = load_workbook(path, data_only=True, read_only=True)

    parts: list[str] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows_data: list[list[str]] = []

        for i, row in enumerate(ws.iter_rows(values_only=True)):
            if i >= max_rows + 1:  # +1 for header
                break
            rows_data.append([_cell_to_str(cell) for cell in row])

        if not rows_data:
            parts.append(f"## Sheet: \"{sheet_name}\" (empty)\n")
            continue

        # Determine dimensions
        n_cols = max(len(r) for r in rows_data) if rows_data else 0
        n_rows = len(rows_data) - 1  # exclude header row

        # Pad rows to equal length
        for row in rows_data:
            while len(row) < n_cols:
                row.append("")

        # Count total rows in sheet (approximate for read_only mode)
        total_rows_note = ""
        if n_rows >= max_rows:
            total_rows_note = f" (showing first {max_rows} rows)"

        parts.append(
            f"## Sheet: \"{sheet_name}\" ({n_rows} rows x {n_cols} columns){total_rows_note}\n"
        )

        # Build markdown table
        if rows_data:
            header = rows_data[0]
            parts.append("| " + " | ".join(header) + " |")
            parts.append("| " + " | ".join("---" for _ in header) + " |")
            for row in rows_data[1:]:
                parts.append("| " + " | ".join(row) + " |")

        parts.append("")  # blank line between sheets

    wb.close()
    return "\n".join(parts).strip()


def _cell_to_str(value: Any) -> str:
    """Convert a cell value to a clean string."""
    if value is None:
        return ""
    if isinstance(value, float):
        # Remove trailing .0 for integers stored as floats
        if value == int(value):
            return str(int(value))
        return f"{value:.2f}"
    return str(value).strip()


# ---------------------------------------------------------------------------
# Word reader
# ---------------------------------------------------------------------------
def read_word(file_path: str) -> str:
    """Read a Word document into structured markdown.

    Extracts headings, paragraphs, bullet lists, and tables.

    Args:
        file_path: Path to .docx file.

    Returns:
        Formatted markdown string.
    """
    from docx import Document
    from docx.oxml.ns import qn

    path = _validate_read_path(file_path)
    doc = Document(path)

    parts: list[str] = []

    for element in doc.element.body:
        tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag

        if tag == "p":
            # It's a paragraph
            para = None
            for p in doc.paragraphs:
                if p._element is element:
                    para = p
                    break
            if para is None:
                continue

            style_name = para.style.name if para.style else ""

            # Headings
            if style_name.startswith("Heading"):
                try:
                    level = int(style_name.split()[-1])
                except (ValueError, IndexError):
                    level = 1
                prefix = "#" * min(level, 6)
                parts.append(f"{prefix} {para.text}")
                parts.append("")
            # List items
            elif _is_list_item(para):
                parts.append(f"- {para.text}")
            # Normal paragraph
            elif para.text.strip():
                parts.append(para.text)
                parts.append("")

        elif tag == "tbl":
            # It's a table
            for table in doc.tables:
                if table._element is element:
                    parts.append(_docx_table_to_md(table))
                    parts.append("")
                    break

    return "\n".join(parts).strip()


def _is_list_item(para: Any) -> bool:
    """Check if a paragraph is a list item by looking for numPr in pPr."""
    from docx.oxml.ns import qn

    pPr = para._element.find(qn("w:pPr"))
    if pPr is not None:
        numPr = pPr.find(qn("w:numPr"))
        if numPr is not None:
            return True
    return False


def _docx_table_to_md(table: Any) -> str:
    """Convert a python-docx table to markdown."""
    rows_data: list[list[str]] = []
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        rows_data.append(cells)

    if not rows_data:
        return ""

    n_cols = max(len(r) for r in rows_data)
    for row in rows_data:
        while len(row) < n_cols:
            row.append("")

    lines: list[str] = []
    header = rows_data[0]
    lines.append("| " + " | ".join(header) + " |")
    lines.append("| " + " | ".join("---" for _ in header) + " |")
    for row in rows_data[1:]:
        lines.append("| " + " | ".join(row) + " |")

    return "\n".join(lines)


# ---------------------------------------------------------------------------
# PowerPoint reader
# ---------------------------------------------------------------------------
def read_presentation(file_path: str) -> str:
    """Read a PowerPoint presentation into structured markdown.

    Extracts slide titles, content, speaker notes, and tables.

    Args:
        file_path: Path to .pptx file.

    Returns:
        Formatted markdown string.
    """
    from pptx import Presentation

    path = _validate_read_path(file_path)
    prs = Presentation(path)

    parts: list[str] = []

    for i, slide in enumerate(prs.slides, 1):
        # Get title
        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text.strip()

        header = f"## Slide {i}"
        if title:
            header += f": \"{title}\""
        parts.append(header)

        # Extract content from shapes
        title_shape_id = slide.shapes.title.shape_id if slide.shapes.title else None
        for shape in slide.shapes:
            # Skip title shape (already captured)
            if shape.shape_id == title_shape_id:
                continue

            # Tables
            if shape.has_table:
                parts.append(_pptx_table_to_md(shape.table))
                continue

            # Text frames (bullets, body text)
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    # Detect bullet level
                    level = para.level if hasattr(para, "level") else 0
                    indent = "  " * level
                    parts.append(f"{indent}- {text}")

        # Speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"\n[Speaker Notes: {notes}]")

        parts.append("")  # blank line between slides

    return "\n".join(parts).strip()


def _pptx_table_to_md(table: Any) -> str:
    """Convert a python-pptx table to markdown."""
    rows_data: list[list[str]] = []
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        rows_data.append(cells)

    if not rows_data:
        return ""

    n_cols = max(len(r) for r in rows_data)
    for row in rows_data:
        while len(row) < n_cols:
            row.append("")

    lines: list[str] = []
    header = rows_data[0]
    lines.append("| " + " | ".join(header) + " |")
    lines.append("| " + " | ".join("---" for _ in header) + " |")
    for row in rows_data[1:]:
        lines.append("| " + " | ".join(row) + " |")

    return "\n".join(lines)


# ---------------------------------------------------------------------------
# PDF reader
# ---------------------------------------------------------------------------
def read_pdf(file_path: str, max_pages: int = 50) -> str:
    """Read a PDF file into structured text.

    Extracts text per page and detects tables.

    Args:
        file_path: Path to .pdf file.
        max_pages: Maximum pages to read.

    Returns:
        Formatted text string.
    """
    import pdfplumber

    path = _validate_read_path(file_path)

    parts: list[str] = []

    with pdfplumber.open(path) as pdf:
        total = len(pdf.pages)
        pages_to_read = min(total, max_pages)

        for i in range(pages_to_read):
            page = pdf.pages[i]
            parts.append(f"## Page {i + 1}")

            # Extract tables first
            tables = page.extract_tables()
            if tables:
                for table in tables:
                    if not table or not table[0]:
                        continue
                    # Build markdown table
                    header = [str(c or "") for c in table[0]]
                    parts.append("| " + " | ".join(header) + " |")
                    parts.append("| " + " | ".join("---" for _ in header) + " |")
                    for row in table[1:]:
                        cells = [str(c or "") for c in row]
                        parts.append("| " + " | ".join(cells) + " |")
                    parts.append("")

            # Extract text (excluding table areas for cleaner output)
            text = page.extract_text()
            if text:
                parts.append(text.strip())

            parts.append("")

        if total > max_pages:
            parts.append(f"\n[Truncated: showing {max_pages} of {total} pages]")

    return "\n".join(parts).strip()


# ---------------------------------------------------------------------------
# CSV reader
# ---------------------------------------------------------------------------
def read_csv(file_path: str, max_rows: int = 200) -> str:
    """Read a CSV or TSV file into a markdown table.

    Auto-detects delimiter (comma, tab, semicolon).

    Args:
        file_path: Path to .csv or .tsv file.
        max_rows: Maximum data rows to read.

    Returns:
        Formatted markdown table string.
    """
    path = _validate_read_path(file_path)

    # Read file content and detect delimiter
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        sample = f.read(4096)
        f.seek(0)

        # Detect delimiter
        delimiter = ","
        if sample.count("\t") > sample.count(","):
            delimiter = "\t"
        elif sample.count(";") > sample.count(","):
            delimiter = ";"

        reader = csv.reader(f, delimiter=delimiter)
        rows: list[list[str]] = []
        for i, row in enumerate(reader):
            if i > max_rows:
                break
            rows.append([cell.strip() for cell in row])

    if not rows:
        return "(Empty CSV file)"

    # Determine column count
    n_cols = max(len(r) for r in rows)
    for row in rows:
        while len(row) < n_cols:
            row.append("")

    # Build markdown table
    parts: list[str] = []
    n_data = len(rows) - 1
    truncated = n_data >= max_rows
    parts.append(f"({n_data} rows x {n_cols} columns)\n")

    header = rows[0]
    parts.append("| " + " | ".join(header) + " |")
    parts.append("| " + " | ".join("---" for _ in header) + " |")
    for row in rows[1:]:
        parts.append("| " + " | ".join(row) + " |")

    if truncated:
        parts.append(f"\n[Truncated: showing first {max_rows} rows]")

    return "\n".join(parts)


# ---------------------------------------------------------------------------
# Plain text reader
# ---------------------------------------------------------------------------
def read_text(file_path: str, max_chars: int = 50_000) -> str:
    """Read a plain text file.

    Supports .txt, .md, .json, .xml, .log, code files, etc.

    Args:
        file_path: Path to text file.
        max_chars: Maximum characters to read.

    Returns:
        File content as string.
    """
    path = _validate_read_path(file_path)

    try:
        with open(path, "r", encoding="utf-8") as f:
            content = f.read(max_chars + 1)
    except UnicodeDecodeError:
        with open(path, "r", encoding="latin-1") as f:
            content = f.read(max_chars + 1)

    truncated = len(content) > max_chars
    if truncated:
        content = content[:max_chars]
        content += f"\n\n[Truncated: showing first {max_chars} characters]"

    return content


# ---------------------------------------------------------------------------
# Main dispatcher
# ---------------------------------------------------------------------------
def parse_file(file_path: str) -> str:
    """Parse any supported file into structured text for LLM context.

    Detects file type by extension and routes to the appropriate reader.

    Args:
        file_path: Path to the file to parse.

    Returns:
        Formatted string with file content, prefixed with filename header.
    """
    ext = os.path.splitext(file_path)[1].lower()
    basename = os.path.basename(file_path)

    try:
        if ext in (".xlsx", ".xls"):
            content = read_excel(file_path)
        elif ext == ".docx":
            content = read_word(file_path)
        elif ext == ".pptx":
            content = read_presentation(file_path)
        elif ext == ".pdf":
            content = read_pdf(file_path)
        elif ext in (".csv", ".tsv"):
            content = read_csv(file_path)
        elif ext in _TEXT_EXTENSIONS:
            content = read_text(file_path)
        else:
            # Try as text file — many unlisted extensions are readable
            try:
                content = read_text(file_path)
            except Exception:
                return f"### {basename}\n(Unsupported file type: {ext})"

        return f"### {basename}\n{content}"

    except FileNotFoundError:
        return f"### {basename}\n(File not found: {file_path})"
    except ValueError as e:
        return f"### {basename}\n(Error: {e})"
    except Exception as e:
        logger.error(f"Error parsing {file_path}: {e}")
        return f"### {basename}\n(Failed to parse: {e})"


def parse_files(file_paths: list[str]) -> str:
    """Parse multiple files and combine their content.

    Args:
        file_paths: List of file paths to parse.

    Returns:
        Combined formatted string with all file contents.
    """
    if not file_paths:
        return ""

    results: list[str] = []
    total_chars = 0

    for path in file_paths:
        result = parse_file(path)
        total_chars += len(result)

        if total_chars > MAX_TOTAL_CHARS:
            remaining = MAX_TOTAL_CHARS - (total_chars - len(result))
            if remaining > 0:
                result = result[:remaining] + "\n\n[Content truncated — file limit reached]"
                results.append(result)
            results.append(
                f"\n[Skipped remaining files — total content exceeded {MAX_TOTAL_CHARS} characters]"
            )
            break

        results.append(result)

    return "\n\n".join(results)


# ---------------------------------------------------------------------------
# Multimodal parsing (text + images)
# ---------------------------------------------------------------------------
def parse_file_multimodal(file_path: str) -> ParseResult:
    """Parse a file into structured text AND extract embedded images.

    Returns a ParseResult with both text content and image data.
    Image extraction is attempted for PPTX, DOCX, and PDF files.
    """
    # Get text content via existing parser
    text = parse_file(file_path)

    # Extract images based on file type
    ext = os.path.splitext(file_path)[1].lower()
    images: list[ImageData] = []

    try:
        real_path = _validate_read_path(file_path)

        if ext == ".pptx":
            images = _extract_pptx_images(real_path)
        elif ext == ".docx":
            images = _extract_docx_images(real_path)
        elif ext == ".pdf":
            images = _extract_pdf_images(real_path)
        elif ext in (".xlsx", ".xls"):
            # No image extraction for Excel — add chart descriptions to text
            chart_info = _describe_excel_charts(real_path)
            if chart_info:
                text += f"\n\n{chart_info}"
    except (FileNotFoundError, ValueError) as e:
        logger.debug("Cannot extract images from %s: %s", file_path, e)
    except Exception as e:
        logger.warning("Image extraction failed for %s: %s", file_path, e)

    return ParseResult(text=text, images=images)


def parse_files_multimodal(file_paths: list[str]) -> ParseResult:
    """Parse multiple files and combine text + images.

    Enforces MAX_IMAGES_TOTAL across all files.

    Returns:
        Combined ParseResult with all text and images.
    """
    if not file_paths:
        return ParseResult(text="")

    texts: list[str] = []
    all_images: list[ImageData] = []
    total_chars = 0

    for path in file_paths:
        result = parse_file_multimodal(path)
        total_chars += len(result.text)

        if total_chars > MAX_TOTAL_CHARS:
            remaining = MAX_TOTAL_CHARS - (total_chars - len(result.text))
            if remaining > 0:
                texts.append(result.text[:remaining] + "\n\n[Content truncated — file limit reached]")
            texts.append(
                f"\n[Skipped remaining files — total content exceeded {MAX_TOTAL_CHARS} characters]"
            )
            break

        texts.append(result.text)

        # Add images up to global limit
        for img in result.images:
            if len(all_images) >= _MAX_IMAGES_TOTAL:
                break
            all_images.append(img)

    return ParseResult(text="\n\n".join(texts), images=all_images)

"""PowerPoint tool — Create professional .pptx files via python-pptx.

Produces consultant-quality presentations with proper bullet formatting,
speaker notes, tables, two-column layouts, and consistent theming.
"""

import os
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from neo.tools.paths import resolve_path, _validate_write_path

# ---------------------------------------------------------------------------
# Theme defaults
# ---------------------------------------------------------------------------
_DEFAULT_THEME = {
    "primary_color": "2B5797",
    "secondary_color": "4472C4",
    "background_color": "FFFFFF",
    "text_color": "333333",
    "font_title": "Calibri",
    "font_body": "Calibri",
    "font_size_title": 28,
    "font_size_subtitle": 18,
    "font_size_body": 18,
    "font_size_bullets": 16,
    "font_size_table": 11,
}

# Slide layout indices (standard PowerPoint template)
_LAYOUT_TITLE = 0  # Title Slide
_LAYOUT_CONTENT = 1  # Title and Content
_LAYOUT_SECTION = 2  # Section Header
_LAYOUT_TWO_CONTENT = 3  # Two Content
_LAYOUT_BLANK = 6  # Blank
_LAYOUT_TITLE_ONLY = 5  # Title Only


def _get_theme(theme: dict | None) -> dict:
    """Merge user theme with defaults."""
    config = dict(_DEFAULT_THEME)
    if theme:
        for key, value in theme.items():
            if key in config and value is not None:
                config[key] = value
    return config


def _hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex string to RGBColor."""
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ---------------------------------------------------------------------------
# Content helpers
# ---------------------------------------------------------------------------
def _add_bullets(text_frame, items: list[str], theme: dict) -> None:
    """Add formatted bullet points to a text frame."""
    text_frame.clear()
    font_name = theme["font_body"]
    font_size = theme["font_size_bullets"]

    for i, item in enumerate(items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = item
        p.font.size = Pt(font_size)
        p.font.name = font_name
        p.font.color.rgb = _hex_to_rgb(theme["text_color"])
        p.space_after = Pt(6)
        p.space_before = Pt(2)
        p.level = 0


def _add_numbered_list(text_frame, items: list[str], theme: dict) -> None:
    """Add a numbered list to a text frame."""
    text_frame.clear()
    font_name = theme["font_body"]
    font_size = theme["font_size_bullets"]

    for i, item in enumerate(items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = f"{i + 1}. {item}"
        p.font.size = Pt(font_size)
        p.font.name = font_name
        p.font.color.rgb = _hex_to_rgb(theme["text_color"])
        p.space_after = Pt(6)
        p.space_before = Pt(2)


def _add_body_text(text_frame, text: str, theme: dict) -> None:
    """Add paragraph text to a text frame."""
    text_frame.clear()
    font_name = theme["font_body"]
    font_size = theme["font_size_body"]

    paragraphs = text.split("\n")
    for i, para in enumerate(paragraphs):
        if not para.strip():
            continue
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = para.strip()
        p.font.size = Pt(font_size)
        p.font.name = font_name
        p.font.color.rgb = _hex_to_rgb(theme["text_color"])
        p.space_after = Pt(8)


def _style_title(shape, title_text: str, theme: dict) -> None:
    """Apply consistent title styling."""
    if shape is None:
        return
    shape.text = title_text
    for paragraph in shape.text_frame.paragraphs:
        paragraph.font.size = Pt(theme["font_size_title"])
        paragraph.font.name = theme["font_title"]
        paragraph.font.bold = True
        paragraph.font.color.rgb = _hex_to_rgb(theme["primary_color"])


def _add_speaker_notes(slide, notes_text: str) -> None:
    """Add speaker notes to a slide."""
    if not notes_text:
        return
    notes_slide = slide.notes_slide
    notes_tf = notes_slide.notes_text_frame
    notes_tf.text = notes_text


def _add_slide_number(slide) -> None:
    """Add a slide number placeholder to the slide.

    Uses the built-in slide number field via XML manipulation.
    """
    from pptx.oxml.ns import qn
    from lxml import etree

    # Create a small text box in the bottom-right for the slide number
    left = Inches(8.5)
    top = Inches(7.0)
    width = Inches(1.0)
    height = Inches(0.4)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT

    # Add slide number field
    run = p.add_run()
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(0x80, 0x80, 0x80)

    fld = run._r.makeelement(qn("a:fld"), {})
    fld.set("type", "slidenum")
    import uuid
    fld.set("{http://schemas.openxmlformats.org/drawingml/2006/main}id", str(uuid.uuid4()).upper())

    # Add run properties and text to the field
    rPr = fld.makeelement(qn("a:rPr"), {"lang": "en-US", "dirty": "0"})
    fld.append(rPr)
    t = fld.makeelement(qn("a:t"), {})
    t.text = "<#>"
    fld.append(t)

    run._r.append(fld)


def _add_table_to_slide(slide, table_data: dict, theme: dict) -> None:
    """Add a professionally styled table to a slide."""
    headers = table_data.get("headers", [])
    rows = table_data.get("rows", [])

    if not headers:
        return

    num_rows = 1 + len(rows)
    num_cols = len(headers)

    # Position the table
    left = Inches(0.8)
    top = Inches(2.0)
    width = Inches(8.4)
    height = Inches(0.4 * num_rows)

    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table

    # Column widths (even distribution)
    col_width = int(width / num_cols)
    for col_idx in range(num_cols):
        table.columns[col_idx].width = col_width

    # Header row styling
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = str(header)

        # Style the cell
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(theme["font_size_table"])
            paragraph.font.name = theme["font_body"]
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            paragraph.alignment = PP_ALIGN.CENTER

        # Header background
        _set_cell_fill(cell, theme["primary_color"])

    # Data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, value in enumerate(row_data):
            if col_idx < num_cols:
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(value) if value is not None else ""

                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(theme["font_size_table"])
                    paragraph.font.name = theme["font_body"]
                    paragraph.font.color.rgb = _hex_to_rgb(theme["text_color"])

                # Alternating row shading
                if row_idx % 2 == 1:
                    _set_cell_fill(cell, "EDF2F9")


def _set_cell_fill(cell, hex_color: str) -> None:
    """Set the background color of a table cell."""
    from pptx.oxml.ns import qn

    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    # Remove existing shading
    for child in tcPr.findall(qn("a:solidFill")):
        tcPr.remove(child)

    solidFill = tcPr.makeelement(qn("a:solidFill"), {})
    srgbClr = solidFill.makeelement(qn("a:srgbClr"), {"val": hex_color.lstrip("#")})
    solidFill.append(srgbClr)
    tcPr.insert(0, solidFill)


def _add_content_to_frame(text_frame, slide_def: dict, theme: dict) -> None:
    """Add the appropriate content to a text frame based on slide definition."""
    if "bullets" in slide_def:
        _add_bullets(text_frame, slide_def["bullets"], theme)
    elif "numbered_list" in slide_def:
        _add_numbered_list(text_frame, slide_def["numbered_list"], theme)
    elif "body" in slide_def:
        _add_body_text(text_frame, slide_def["body"], theme)
    elif "content" in slide_def:
        # Backward compatibility: treat 'content' as body text
        _add_body_text(text_frame, slide_def["content"], theme)


# ---------------------------------------------------------------------------
# Main function
# ---------------------------------------------------------------------------
def create_presentation(
    title: str,
    slides: list[dict] | None = None,
    output_path: str | None = None,
    theme: dict | None = None,
) -> str:
    """Create a professional PowerPoint presentation.

    Args:
        title: Filename (without extension) or full path.
        slides: List of slide dicts. Each slide supports:
            - title (str): Slide heading
            - layout (str): "title", "content", "section_header", "two_column",
                           "blank", "title_only" (default: auto-detect)
            - bullets (list[str]): Bullet point list
            - numbered_list (list[str]): Numbered list
            - body (str): Paragraph text
            - content (str): Deprecated, treated as body text
            - table ({headers: [...], rows: [[...], ...]}): Data table
            - speaker_notes (str): Presenter notes
            - left_content (dict): For two_column layout — {bullets/body}
            - right_content (dict): For two_column layout — {bullets/body}
        output_path: Directory where the file should be saved.
        theme: Color and font settings:
            - primary_color (str): Hex color for titles/accents (default "2B5797")
            - secondary_color (str): Hex for secondary elements
            - text_color (str): Hex for body text (default "333333")
            - font_title (str): Title font (default "Calibri")
            - font_body (str): Body font (default "Calibri")
            - font_size_title (int): Title size in pt (default 28)
            - font_size_bullets (int): Bullet size in pt (default 16)

    Returns:
        Absolute path to the created .pptx file.
    """
    prs = Presentation()
    th = _get_theme(theme)

    if not slides:
        # Single title slide
        layout = prs.slide_layouts[_LAYOUT_TITLE]
        slide = prs.slides.add_slide(layout)
        _style_title(slide.shapes.title, os.path.basename(title), th)
        if len(slide.placeholders) > 1 and slide.placeholders[1]:
            slide.placeholders[1].text = "Created by Neo"
    else:
        for i, slide_def in enumerate(slides):
            slide_title = slide_def.get("title", f"Slide {i + 1}")
            layout_name = slide_def.get("layout", "")

            # Auto-detect layout if not specified
            if not layout_name:
                if i == 0 and not any(k in slide_def for k in ("bullets", "body", "table", "numbered_list")):
                    layout_name = "title"
                elif "left_content" in slide_def and "right_content" in slide_def:
                    layout_name = "two_column"
                elif "table" in slide_def:
                    layout_name = "title_only"
                else:
                    layout_name = "content"

            # Map layout name to index
            layout_map = {
                "title": _LAYOUT_TITLE,
                "content": _LAYOUT_CONTENT,
                "section_header": _LAYOUT_SECTION,
                "two_column": _LAYOUT_TWO_CONTENT,
                "blank": _LAYOUT_BLANK,
                "title_only": _LAYOUT_TITLE_ONLY,
            }
            layout_idx = layout_map.get(layout_name, _LAYOUT_CONTENT)

            # Safely get layout (fall back to content if index doesn't exist)
            try:
                layout = prs.slide_layouts[layout_idx]
            except IndexError:
                layout = prs.slide_layouts[_LAYOUT_CONTENT]

            slide = prs.slides.add_slide(layout)

            # --- Title ---
            if layout_name == "title":
                _style_title(slide.shapes.title, slide_title, th)
                # Subtitle
                subtitle_text = slide_def.get("body", slide_def.get("content", ""))
                if subtitle_text and len(slide.placeholders) > 1:
                    ph = slide.placeholders[1]
                    ph.text = subtitle_text
                    for p in ph.text_frame.paragraphs:
                        p.font.size = Pt(th["font_size_subtitle"])
                        p.font.name = th["font_body"]
                        p.font.color.rgb = _hex_to_rgb(th["text_color"])

            elif layout_name == "section_header":
                _style_title(slide.shapes.title, slide_title, th)
                if len(slide.placeholders) > 1 and slide_def.get("body"):
                    ph = slide.placeholders[1]
                    ph.text = slide_def["body"]

            elif layout_name == "two_column":
                _style_title(slide.shapes.title, slide_title, th)
                left = slide_def.get("left_content", {})
                right = slide_def.get("right_content", {})

                # Left content (placeholder index 1)
                if len(slide.placeholders) > 1 and left:
                    _add_content_to_frame(slide.placeholders[1].text_frame, left, th)

                # Right content (placeholder index 2)
                if len(slide.placeholders) > 2 and right:
                    _add_content_to_frame(slide.placeholders[2].text_frame, right, th)

            elif layout_name == "title_only":
                _style_title(slide.shapes.title, slide_title, th)

            else:
                # Standard content slide
                _style_title(slide.shapes.title, slide_title, th)

                if len(slide.placeholders) > 1:
                    body_ph = slide.placeholders[1]
                    _add_content_to_frame(body_ph.text_frame, slide_def, th)

            # --- Table (added as shape on any layout) ---
            if "table" in slide_def:
                _add_table_to_slide(slide, slide_def["table"], th)

            # --- Speaker notes ---
            notes = slide_def.get("speaker_notes", "")
            if notes:
                _add_speaker_notes(slide, notes)

            # --- Slide number (skip title slide) ---
            if layout_name != "title":
                _add_slide_number(slide)

    # Save
    file_path = resolve_path(title, ".pptx", output_dir=output_path)
    os.makedirs(os.path.dirname(file_path), exist_ok=True)
    prs.save(file_path)
    return file_path


# ---------------------------------------------------------------------------
# Collaborative editing
# ---------------------------------------------------------------------------
def edit_presentation(
    file_path: str,
    operations: list[dict] | None = None,
    save_as: str | None = None,
) -> str:
    """Edit an existing PowerPoint presentation — update/add/delete slides.

    Args:
        file_path: Path to the existing .pptx file.
        operations: List of editing operations. Each dict has a "type" key:
            - update_slide: {type, slide (1-based), title, bullets, body, notes}
              Updates an existing slide's content. Only provided fields are changed.
            - add_slide: {type, title, bullets, body, speaker_notes}
              Appends a new content slide at the end.
            - delete_slide: {type, slide (1-based)}
              Removes a slide by number.
            - update_notes: {type, slide (1-based), notes}
              Sets or replaces speaker notes on a slide.
        save_as: Optional alternate save path. If omitted, overwrites original.

    Returns:
        Absolute path to the saved file with a summary of changes.
    """
    from neo.tools.file_reader import _validate_read_path

    real_path = _validate_read_path(file_path)
    prs = Presentation(real_path)
    th = _get_theme(None)

    changes: list[str] = []

    if operations:
        # Collect delete indices to process after other ops (in reverse)
        deletes: list[int] = []

        for op in operations:
            op_type = op.get("type", "")

            if op_type == "update_slide":
                slide_num = op.get("slide", 0)
                if 1 <= slide_num <= len(prs.slides):
                    slide = prs.slides[slide_num - 1]

                    # Update title
                    new_title = op.get("title")
                    if new_title is not None and slide.shapes.title:
                        _style_title(slide.shapes.title, new_title, th)

                    # Update body content (find the first non-title text placeholder)
                    if any(k in op for k in ("bullets", "body")):
                        title_id = slide.shapes.title.shape_id if slide.shapes.title else None
                        for shape in slide.shapes:
                            if shape.shape_id == title_id:
                                continue
                            if shape.has_text_frame:
                                if "bullets" in op:
                                    _add_bullets(shape.text_frame, op["bullets"], th)
                                elif "body" in op:
                                    _add_body_text(shape.text_frame, op["body"], th)
                                break

                    # Update notes
                    if "notes" in op:
                        _add_speaker_notes(slide, op["notes"])

                    changes.append(f"Updated slide {slide_num}")

            elif op_type == "add_slide":
                # Add a content slide at the end
                try:
                    layout = prs.slide_layouts[_LAYOUT_CONTENT]
                except IndexError:
                    layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(layout)

                title_text = op.get("title", "")
                if title_text and slide.shapes.title:
                    _style_title(slide.shapes.title, title_text, th)

                # Add content to body placeholder
                if len(slide.placeholders) > 1:
                    tf = slide.placeholders[1].text_frame
                    if "bullets" in op:
                        _add_bullets(tf, op["bullets"], th)
                    elif "body" in op:
                        _add_body_text(tf, op["body"], th)

                notes = op.get("speaker_notes", "")
                if notes:
                    _add_speaker_notes(slide, notes)

                _add_slide_number(slide)
                changes.append(f"Added slide '{title_text}'")

            elif op_type == "delete_slide":
                slide_num = op.get("slide", 0)
                if 1 <= slide_num <= len(prs.slides):
                    deletes.append(slide_num)

            elif op_type == "update_notes":
                slide_num = op.get("slide", 0)
                notes = op.get("notes", "")
                if 1 <= slide_num <= len(prs.slides):
                    _add_speaker_notes(prs.slides[slide_num - 1], notes)
                    changes.append(f"Updated notes on slide {slide_num}")

        # Process deletes in reverse order
        for slide_num in sorted(deletes, reverse=True):
            if 1 <= slide_num <= len(prs.slides):
                rId = prs.slides._sldIdLst[slide_num - 1].get(
                    "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
                )
                prs.part.drop_rel(rId)
                slide_elem = prs.slides._sldIdLst[slide_num - 1]
                prs.slides._sldIdLst.remove(slide_elem)
                changes.append(f"Deleted slide {slide_num}")

    # --- Save ---
    if save_as:
        out_path = os.path.expanduser(save_as)
        if not out_path.endswith(".pptx"):
            out_path += ".pptx"
        _validate_write_path(out_path)
    else:
        out_path = real_path

    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    prs.save(out_path)

    summary = "; ".join(changes) if changes else "No changes applied"
    return f"Saved {out_path} ({summary})"

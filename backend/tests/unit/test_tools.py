"""Tests for all tool modules — Excel, PowerPoint, Word, Obsidian, Files."""

import os
import tempfile

import pytest
from openpyxl import load_workbook
from pptx import Presentation

from neo.tools.excel import create_workbook
from neo.tools.files import manage_file
from neo.tools.obsidian import append_to_note, create_note
from neo.tools.powerpoint import create_presentation
from neo.tools.word import create_document


@pytest.fixture
def tmp_dir():
    with tempfile.TemporaryDirectory() as d:
        yield d


# ============================================
# EXCEL
# ============================================


class TestExcel:
    def test_create_empty_workbook(self, tmp_dir):
        path = create_workbook(title=os.path.join(tmp_dir, "Test"))
        assert os.path.exists(path)
        assert path.endswith(".xlsx")

    def test_create_with_sheets_and_data(self, tmp_dir):
        sheets = [
            {
                "name": "Sales",
                "headers": ["Month", "Revenue", "Profit"],
                "rows": [
                    ["January", 10000, 3000],
                    ["February", 12000, 4000],
                ],
            },
            {
                "name": "Summary",
                "headers": ["Metric", "Value"],
                "rows": [["Total Revenue", 22000]],
            },
        ]
        path = create_workbook(title=os.path.join(tmp_dir, "Report"), sheets=sheets)
        assert os.path.exists(path)

        # Verify structure
        wb = load_workbook(path)
        assert len(wb.sheetnames) == 2
        assert "Sales" in wb.sheetnames
        assert "Summary" in wb.sheetnames

        ws = wb["Sales"]
        assert ws.cell(1, 1).value == "Month"
        assert ws.cell(1, 2).value == "Revenue"
        assert ws.cell(2, 1).value == "January"
        assert ws.cell(2, 2).value == 10000
        assert ws.cell(3, 2).value == 12000
        wb.close()

    def test_headers_are_styled(self, tmp_dir):
        sheets = [{"name": "Data", "headers": ["A", "B"], "rows": []}]
        path = create_workbook(title=os.path.join(tmp_dir, "Styled"), sheets=sheets)
        wb = load_workbook(path)
        ws = wb["Data"]
        cell = ws.cell(1, 1)
        assert cell.font.bold is True
        wb.close()

    def test_freeze_panes_on_header(self, tmp_dir):
        sheets = [{"name": "Data", "headers": ["A"], "rows": []}]
        path = create_workbook(title=os.path.join(tmp_dir, "Freeze"), sheets=sheets)
        wb = load_workbook(path)
        ws = wb["Data"]
        assert ws.freeze_panes == "A2"
        wb.close()


# ============================================
# POWERPOINT
# ============================================


class TestPowerPoint:
    def test_create_empty_presentation(self, tmp_dir):
        path = create_presentation(title=os.path.join(tmp_dir, "Deck"))
        assert os.path.exists(path)
        assert path.endswith(".pptx")

    def test_create_with_slides(self, tmp_dir):
        slides = [
            {"title": "Welcome", "content": "Introduction to the project"},
            {"title": "Overview", "content": "Key metrics and findings"},
            {"title": "Conclusion", "content": "Next steps"},
        ]
        path = create_presentation(title=os.path.join(tmp_dir, "Report"), slides=slides)
        assert os.path.exists(path)

        prs = Presentation(path)
        assert len(prs.slides) == 3
        assert prs.slides[0].shapes.title.text == "Welcome"

    def test_single_slide_has_title(self, tmp_dir):
        path = create_presentation(title=os.path.join(tmp_dir, "Single"))
        prs = Presentation(path)
        assert len(prs.slides) == 1
        assert prs.slides[0].shapes.title.text == "Single"


# ============================================
# WORD
# ============================================


class TestWord:
    def test_create_empty_document(self, tmp_dir):
        path = create_document(title=os.path.join(tmp_dir, "Doc"))
        assert os.path.exists(path)
        assert path.endswith(".docx")

    def test_create_with_content(self, tmp_dir):
        content = "# Introduction\nThis is the first section.\n## Details\n- Item one\n- Item two"
        path = create_document(title=os.path.join(tmp_dir, "Report"), content=content)
        assert os.path.exists(path)

        from docx import Document

        doc = Document(path)
        # Title + headings + paragraphs + list items
        assert len(doc.paragraphs) >= 4

    def test_headings_converted(self, tmp_dir):
        content = "# Heading 1\n## Heading 2\n### Heading 3"
        path = create_document(title=os.path.join(tmp_dir, "Headings"), content=content)

        from docx import Document

        doc = Document(path)
        heading_styles = [p.style.name for p in doc.paragraphs if "Heading" in p.style.name]
        assert len(heading_styles) >= 3


# ============================================
# OBSIDIAN
# ============================================


class TestObsidian:
    def test_create_note_basic(self, tmp_dir, monkeypatch):
        monkeypatch.setattr("neo.tools.obsidian._get_default_vault", lambda: tmp_dir)
        path = create_note(title="Test Note", content="Hello world")
        assert os.path.exists(path)
        assert path.endswith(".md")

        with open(path) as f:
            text = f.read()
        assert "---" in text  # frontmatter
        assert "title: Test Note" in text
        assert "# Test Note" in text
        assert "Hello world" in text

    def test_create_note_with_tags(self, tmp_dir, monkeypatch):
        monkeypatch.setattr("neo.tools.obsidian._get_default_vault", lambda: tmp_dir)
        path = create_note(title="Tagged", tags=["project", "neo"])

        with open(path) as f:
            text = f.read()
        assert "tags: [project, neo]" in text

    def test_create_note_with_backlinks(self, tmp_dir, monkeypatch):
        monkeypatch.setattr("neo.tools.obsidian._get_default_vault", lambda: tmp_dir)
        path = create_note(title="Linked", links=["Other Note", "Reference"])

        with open(path) as f:
            text = f.read()
        assert "[[Other Note]]" in text
        assert "[[Reference]]" in text

    def test_frontmatter_has_date(self, tmp_dir, monkeypatch):
        monkeypatch.setattr("neo.tools.obsidian._get_default_vault", lambda: tmp_dir)
        path = create_note(title="Dated")

        with open(path) as f:
            text = f.read()
        assert "date:" in text
        assert "created_by: neo" in text

    def test_append_to_note(self, tmp_dir, monkeypatch):
        monkeypatch.setattr("neo.tools.obsidian._get_default_vault", lambda: tmp_dir)
        path = create_note(title="Append Test", content="Initial")
        append_to_note(path, "Added later")

        with open(path) as f:
            text = f.read()
        assert "Initial" in text
        assert "Added later" in text

    def test_vault_override_used(self, tmp_dir, monkeypatch):
        from neo.tools.obsidian import set_vault_path
        set_vault_path(tmp_dir)
        try:
            path = create_note(title="Override Test", content="Works")
            assert path.startswith(tmp_dir)
            assert os.path.exists(path)
        finally:
            set_vault_path(None)

    def test_convert_windows_path_drive(self, monkeypatch):
        from neo.tools.obsidian import _convert_windows_path
        monkeypatch.setattr("neo.tools.obsidian.platform.system", lambda: "Linux")
        assert _convert_windows_path(r"G:\Meu Drive\notes") == "/mnt/g/Meu Drive/notes"
        assert _convert_windows_path(r"C:\Users\andre\vault") == "/mnt/c/Users/andre/vault"

    def test_convert_windows_path_forward_slash(self, monkeypatch):
        from neo.tools.obsidian import _convert_windows_path
        monkeypatch.setattr("neo.tools.obsidian.platform.system", lambda: "Linux")
        assert _convert_windows_path("G:/Meu Drive/notes") == "/mnt/g/Meu Drive/notes"

    def test_convert_windows_path_noop_linux(self):
        from neo.tools.obsidian import _convert_windows_path
        assert _convert_windows_path("/home/user/vault") == "/home/user/vault"

    def test_convert_windows_path_noop_on_windows(self, monkeypatch):
        from neo.tools.obsidian import _convert_windows_path
        monkeypatch.setattr("neo.tools.obsidian.platform.system", lambda: "Windows")
        # On actual Windows, don't convert
        assert _convert_windows_path(r"G:\vault") == r"G:\vault"

    def test_wsl_to_windows(self):
        from neo.tools.obsidian import _wsl_to_windows
        assert _wsl_to_windows("/mnt/g/Meu Drive/notes") == r"G:\Meu Drive\notes"
        assert _wsl_to_windows("/mnt/c/Users/andre") == r"C:\Users\andre"
        # Non-matching paths returned as-is
        assert _wsl_to_windows("/home/user/vault") == "/home/user/vault"

    def test_needs_windows_io_unmounted(self, monkeypatch):
        from neo.tools.obsidian import _needs_windows_io
        monkeypatch.setattr("os.path.isdir", lambda p: p == "/mnt/c")
        assert _needs_windows_io("/mnt/g/Meu Drive/vault") is True
        assert _needs_windows_io("/mnt/c/Users/vault") is False
        assert _needs_windows_io("/home/user/vault") is False

    def test_write_file_powershell_fallback(self, monkeypatch):
        """When drive is not mounted, _write_file should call PowerShell."""
        from unittest.mock import MagicMock, patch

        from neo.tools.obsidian import _write_file

        monkeypatch.setattr("neo.tools.obsidian._needs_windows_io", lambda p: True)
        mock_result = MagicMock(returncode=0, stderr="")
        with patch("neo.tools.obsidian.subprocess.run", return_value=mock_result) as mock_run:
            _write_file("/mnt/g/Vault/Note.md", "Hello")
            mock_run.assert_called_once()
            call_args = mock_run.call_args
            assert call_args[0][0][0] == "powershell.exe"
            assert call_args[1]["input"] == "Hello"

    def test_append_file_powershell_fallback(self, monkeypatch):
        """When drive is not mounted, _append_file should call PowerShell."""
        from unittest.mock import MagicMock, patch

        from neo.tools.obsidian import _append_file

        monkeypatch.setattr("neo.tools.obsidian._needs_windows_io", lambda p: True)
        mock_result = MagicMock(returncode=0, stderr="")
        with patch("neo.tools.obsidian.subprocess.run", return_value=mock_result) as mock_run:
            _append_file("/mnt/g/Vault/Note.md", "Extra content")
            mock_run.assert_called_once()
            call_args = mock_run.call_args
            assert call_args[0][0][0] == "powershell.exe"
            assert "Extra content" in call_args[1]["input"]

    def test_write_file_powershell_error(self, monkeypatch):
        """PowerShell failure should raise OSError."""
        from unittest.mock import MagicMock, patch

        from neo.tools.obsidian import _write_file

        monkeypatch.setattr("neo.tools.obsidian._needs_windows_io", lambda p: True)
        mock_result = MagicMock(returncode=1, stderr="Access denied")
        with patch("neo.tools.obsidian.subprocess.run", return_value=mock_result):
            with pytest.raises(OSError, match="PowerShell write failed"):
                _write_file("/mnt/g/Vault/Note.md", "Hello")

    def test_create_note_windows_io(self, monkeypatch):
        """create_note should use PowerShell when drive isn't mounted."""
        from unittest.mock import MagicMock, patch

        monkeypatch.setattr("neo.tools.obsidian._get_default_vault", lambda: "/mnt/g/Vault")
        monkeypatch.setattr("neo.tools.obsidian._needs_windows_io", lambda p: True)
        mock_result = MagicMock(returncode=0, stderr="")
        with patch("neo.tools.obsidian.subprocess.run", return_value=mock_result) as mock_run:
            path = create_note(title="Test", content="Hello")
            assert path == "/mnt/g/Vault/Test.md"
            mock_run.assert_called_once()
            # Content should contain the frontmatter and body
            written = mock_run.call_args[1]["input"]
            assert "title: Test" in written
            assert "Hello" in written


# ============================================
# FILE SYSTEM
# ============================================


class TestFiles:
    def test_move_file(self, tmp_dir):
        src = os.path.join(tmp_dir, "source.txt")
        dst = os.path.join(tmp_dir, "dest.txt")
        with open(src, "w") as f:
            f.write("content")

        result = manage_file("move", src, dst)
        assert "Moved" in result
        assert not os.path.exists(src)
        assert os.path.exists(dst)

    def test_copy_file(self, tmp_dir):
        src = os.path.join(tmp_dir, "source.txt")
        dst = os.path.join(tmp_dir, "copy.txt")
        with open(src, "w") as f:
            f.write("content")

        result = manage_file("copy", src, dst)
        assert "Copied" in result
        assert os.path.exists(src)  # original still there
        assert os.path.exists(dst)

    def test_rename_file(self, tmp_dir):
        src = os.path.join(tmp_dir, "old.txt")
        with open(src, "w") as f:
            f.write("content")

        result = manage_file("rename", src, os.path.join(tmp_dir, "new.txt"))
        assert "Renamed" in result
        assert not os.path.exists(src)
        assert os.path.exists(os.path.join(tmp_dir, "new.txt"))

    def test_delete_file(self, tmp_dir):
        src = os.path.join(tmp_dir, "to_delete.txt")
        with open(src, "w") as f:
            f.write("content")

        result = manage_file("delete", src)
        assert "Deleted" in result
        assert not os.path.exists(src)

    def test_source_not_found(self):
        with pytest.raises(ValueError, match="does not exist"):
            manage_file("move", "/nonexistent/file.txt", "/tmp/dst.txt")

    def test_delete_directory_refused(self, tmp_dir):
        with pytest.raises(ValueError, match="[Rr]efusing"):
            manage_file("delete", tmp_dir)

    def test_unknown_action(self, tmp_dir):
        src = os.path.join(tmp_dir, "file.txt")
        with open(src, "w") as f:
            f.write("x")
        with pytest.raises(ValueError, match="[Uu]nknown action"):
            manage_file("explode", src)

    def test_safety_check_system_dir(self):
        with pytest.raises(ValueError, match="protected"):
            manage_file("delete", "/bin/bash")


# ============================================
# PATH RESOLUTION (security-critical)
# ============================================


class TestPaths:
    def test_resolve_path_relative_title(self, monkeypatch):
        from neo.tools.paths import resolve_path

        monkeypatch.setenv("DEFAULT_SAVE_DIR", "/tmp/neo_test")
        result = resolve_path("Report", ".xlsx")
        assert result == "/tmp/neo_test/Report.xlsx"

    def test_resolve_path_absolute_valid(self, tmp_dir):
        from neo.tools.paths import resolve_path

        path = os.path.join(tmp_dir, "file")
        result = resolve_path(path, ".xlsx")
        assert result == path + ".xlsx"

    def test_resolve_path_blocks_system_dir(self):
        from neo.tools.paths import resolve_path

        with pytest.raises(ValueError, match="protected"):
            resolve_path("/etc/passwd", ".xlsx")

    def test_resolve_path_blocks_sensitive_home(self):
        from neo.tools.paths import resolve_path

        with pytest.raises(ValueError, match="sensitive"):
            resolve_path(os.path.expanduser("~/.ssh/key"), ".xlsx")

    def test_resolve_path_sanitizes_special_chars(self, monkeypatch):
        from neo.tools.paths import resolve_path

        monkeypatch.setenv("DEFAULT_SAVE_DIR", "/tmp/neo_test")
        result = resolve_path("../../../etc/passwd", ".xlsx")
        # Should be sanitized to safe filename, not a traversal
        assert "/tmp/neo_test/" in result
        assert ".." not in result

    def test_resolve_path_tilde_expansion(self, tmp_dir, monkeypatch):
        from neo.tools.paths import resolve_path

        # Tilde should be expanded
        result = resolve_path("~/Documents/Neo/test", ".xlsx")
        assert "~" not in result
        assert os.path.isabs(result)

    def test_resolve_path_with_output_dir(self, tmp_dir):
        from neo.tools.paths import resolve_path

        result = resolve_path("Report", ".xlsx", output_dir=tmp_dir)
        assert result == os.path.join(tmp_dir, "Report.xlsx")

    def test_resolve_path_output_dir_blocks_system(self):
        from neo.tools.paths import resolve_path

        with pytest.raises(ValueError, match="protected"):
            resolve_path("Report", ".xlsx", output_dir="/etc")

    def test_resolve_path_output_dir_blocks_sensitive(self):
        from neo.tools.paths import resolve_path

        with pytest.raises(ValueError, match="sensitive"):
            resolve_path("key", ".xlsx", output_dir=os.path.expanduser("~/.ssh"))

    def test_arbitrary_user_dir_allowed(self, tmp_dir):
        """Paths outside ~/Documents/Neo are now allowed (relaxed validation)."""
        from neo.tools.paths import resolve_path

        # Any user-writable directory should work
        result = resolve_path(os.path.join(tmp_dir, "subdir", "file"), ".docx")
        assert result.endswith(".docx")
        assert tmp_dir in result


class TestOutputPath:
    """Test output_path parameter on creation tools."""

    def test_word_output_path(self, tmp_dir):
        path = create_document(title="Doc", output_path=tmp_dir)
        assert os.path.exists(path)
        assert path.startswith(tmp_dir)
        assert path.endswith(".docx")

    def test_excel_output_path(self, tmp_dir):
        path = create_workbook(title="Sheet", output_path=tmp_dir)
        assert os.path.exists(path)
        assert path.startswith(tmp_dir)
        assert path.endswith(".xlsx")

    def test_powerpoint_output_path(self, tmp_dir):
        path = create_presentation(title="Deck", output_path=tmp_dir)
        assert os.path.exists(path)
        assert path.startswith(tmp_dir)
        assert path.endswith(".pptx")

    def test_output_path_blocks_system_dir(self):
        with pytest.raises(ValueError, match="protected"):
            create_document(title="Doc", output_path="/etc")

    def test_output_path_blocks_sensitive_dir(self):
        with pytest.raises(ValueError, match="sensitive"):
            create_document(title="Doc", output_path=os.path.expanduser("~/.ssh"))
